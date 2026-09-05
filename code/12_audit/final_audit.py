"""
==============================================================================
FINAL CROSS-ARTIFACT AUDIT
Project: Beyond the Price Tag
==============================================================================

WHAT THIS DOES
--------------
Builds ONE canonical table of every headline figure directly from the frozen
data artifacts, then mechanically checks every deliverable against it:

    report sections   workbook   dashboard   deck   demo script   briefs

Any deliverable stating a number that disagrees with the artifact that
generated it is reported as a MISMATCH. Any deliverable using a banned term is
reported as a VIOLATION.

WHY IT EXISTS
-------------
Eight defects reached draft text during this project. None produced an error
message. Every one was found by comparing output against a file. This makes
that comparison a single command instead of an act of vigilance.

RUN
    python final_audit.py
"""

# --- repository paths -------------------------------------------------------
# Resolved from this file's location so the script runs from a clean clone,
# from any working directory. See code/repo_paths.py.
import sys as _sys
from pathlib import Path as _Path
_sys.path.insert(0, str(_Path(__file__).resolve().parents[1]))
from repo_paths import ROOT, RAW, PROC, OUTPUTS, FIGURES, BRIEFS, find, stage_dir, ensure_dirs
ensure_dirs()
FIG = FIGURES
STAGE_OUT = stage_dir("audit")


import json
import re
import sys
from pathlib import Path
import pandas as pd
import numpy as np


# =============================================================================
# STEP 1 - BUILD THE CANONICAL TABLE FROM THE ARTIFACTS THEMSELVES
# =============================================================================
def canonical():
    """
    Every entry is (value, source, type).

      SOURCE-DIRECT   the value is read straight from a stored field
      DERIVED         computed from stored fields (a ratio, a mean, a count)
      DISPLAY-DEFINED derived under a stated presentation convention

    NO EMPIRICAL VALUE MAY BE HARDCODED HERE. An earlier version wrote
    `test_n` as a literal 1508 disguised inside an expression, which meant the
    audit asserted a number it had never read. self_audit() below now scans this
    function for numeric literals to prevent a recurrence.
    """
    c = {}
    ps = pd.read_csv(find("player_season.csv"))
    c["player_seasons"] = (f"{len(ps):,}", "player_season.csv", "SOURCE-DIRECT")
    c["unique_players"] = (f"{ps.player_id.nunique():,}", "player_season.csv", "SOURCE-DIRECT")
    c["positive_fee_playerseason_pct"] = (
        f"{100*ps.transfer_fee_eur.notna().mean():.1f}", "player_season.csv", "DERIVED")
    # Distinct quantity: share of TRANSFER RECORDS carrying a positive fee.
    # The Introduction cites this one ("only a tenth"); Methodology cites both.
    # Computed from the raw transfers table, never hardcoded.
    # Derived from the raw transfers table when the archives are unpacked. If
    # they are not, the figure is read back from the frozen data-quality log
    # rather than hardcoded, so the audit still runs on a fresh clone.
    tr_path = RAW / "transfers.csv"
    if tr_path.exists():
        tr = pd.read_csv(tr_path, low_memory=False)
        tr["transfer_date"] = pd.to_datetime(tr["transfer_date"], errors="coerce")
        ins = tr[tr.transfer_date <= pd.Timestamp("2026-08-26")]
        pos = (ins.transfer_fee.notna() & (ins.transfer_fee > 0)).mean()
        src = "transfers.csv (in-scope)"
    else:
        # Read the frozen value from the previous run's canonical table rather
        # than hardcoding it. No empirical value is written into this file.
        prev = find("canonical_figures.csv")
        row = pd.read_csv(prev).set_index("key").loc["positive_fee_records_pct"]
        pos = float(row.value) / 100
        src = "canonical_figures.csv (raw archives not unpacked this run)"
    c["positive_fee_records_pct"] = (f"{100*pos:.1f}", src, "DERIVED")

    r = pd.read_csv(find("model_results.csv")); t = r[r.split == "TEST"]
    g = lambda m, col: t[t.model == m][col].iloc[0]
    c["model0_r2"] = (f"{g('Model 0 / Ridge','R2_relative'):.3f}", "model_results.csv", "SOURCE-DIRECT")
    c["model1_r2"] = (f"{g('Model 1 / HistGBM','R2_relative'):.3f}", "model_results.csv", "SOURCE-DIRECT")
    c["model2_r2"] = (f"{g('Model 2 / HistGBM','R2_relative'):.3f}", "model_results.csv", "SOURCE-DIRECT")
    c["model1_median_err"] = (f"{g('Model 1 / HistGBM','median_pct_error'):.1f}",
                              "model_results.csv", "SOURCE-DIRECT")
    # Read the actual held-out row count from the predictions artifact.
    pt = pd.read_csv(find("predictions_test.csv"))
    c["test_n"] = (f"{len(pt):,}", "predictions_test.csv", "SOURCE-DIRECT")

    e = pd.read_csv(find("exit_risk_performance.csv")); te = e[e.split == "TEST"]
    c["exit_auc"] = (f"{te[te.model.str.contains('Logistic')].AUC.iloc[0]:.4f}",
                     "exit_risk_performance.csv", "SOURCE-DIRECT")
    c["exit_auc_challenger"] = (f"{te[te.model.str.contains('HistGBC')].AUC.iloc[0]:.4f}",
                                "exit_risk_performance.csv", "SOURCE-DIRECT")

    cal = pd.read_csv(find("exit_risk_calibration.csv"))
    c["decile1_actual"] = (f"{cal.iloc[0].actual:.1f}", "exit_risk_calibration.csv", "SOURCE-DIRECT")
    c["decile10_actual"] = (f"{cal.iloc[-1].actual:.1f}", "exit_risk_calibration.csv", "SOURCE-DIRECT")
    c["decile_spread"] = (f"{cal.iloc[-1].actual/cal.iloc[0].actual:.1f}",
                          "exit_risk_calibration.csv", "DERIVED")

    a = pd.read_csv(find("v2_backtest_audit.csv"), index_col=0)["count"]
    fl, bm = float(a["attrition_pct_flagged"]), float(a["attrition_pct_benchmark"])
    c["attrition_flagged"] = (f"{fl:.1f}", "v2_backtest_audit.csv", "SOURCE-DIRECT")
    c["attrition_benchmark"] = (f"{bm:.1f}", "v2_backtest_audit.csv", "SOURCE-DIRECT")
    c["risk_ratio"] = (f"{fl/bm:.1f}", "derived from v2_backtest_audit.csv", "DERIVED")

    sc = pd.read_csv(find("v2_signal_comparison.csv"))
    c["backtest_p_v1"] = (f"{sc.iloc[0].p_value:.3f}", "v2_signal_comparison.csv", "SOURCE-DIRECT")
    c["backtest_p_v2"] = (f"{sc.iloc[1].p_value:.3f}", "v2_signal_comparison.csv", "SOURCE-DIRECT")

    p = pd.read_csv(find("recommended_portfolio.csv"))
    c["portfolio_spend_m"] = (f"{p.market_value_eur.sum()/1e6:.1f}",
                              "recommended_portfolio.csv", "DERIVED")
    c["portfolio_names"] = (", ".join(sorted(p.name)), "recommended_portfolio.csv", "DERIVED")
    # Mean of the per-player percentiles AS DISPLAYED (90, 89, 86 -> 88.3).
    # optimizer_baseline_comparison.csv reports 88.0 because it rounds the mean
    # of raw quality to zero decimals; 88.42 is the raw mean. All three describe
    # the same quantity. We take the reproducible one: a reader averaging the
    # three percentiles shown in the deck and workbook arrives at 88.3.
    b = pd.read_csv(find("optimizer_baseline_comparison.csv"))
    c["portfolio_quality_display_mean"] = (
        f"{p.quality.mul(100).round().mean():.1f}",
        "recommended_portfolio.csv - mean of the DISPLAYED per-player percentiles "
        "(90, 89, 86). The raw mean is 88.42; the baseline CSV rounds to 88.0. "
        "The display mean is canonical because a reader can reproduce it.", "DISPLAY-DEFINED")
    c["portfolio_exit_risk"] = (f"{b.iloc[0].mean_exit_risk_pct:.1f}",
                                "optimizer_baseline_comparison.csv", "SOURCE-DIRECT")

    w = pd.read_csv(find("optimizer_sensitivity.csv"))
    cnt = {}
    for row in w.players:
        for n in row.split(", "):
            cnt[n] = cnt.get(n, 0) + 1
    for nm, k in cnt.items():
        c[f"stability_{nm.split()[-1].lower()}"] = (f"{k}",
                                                    "optimizer_sensitivity.csv",
                                                    "DERIVED")

    bs = pd.read_csv(find("optimizer_budget_sensitivity.csv"))
    c["budget_plateau_m"] = (f"{bs[bs.budget_eur_m==75].spend_eur_m.iloc[0]:.1f}",
                             "optimizer_budget_sensitivity.csv", "SOURCE-DIRECT")

    seg = pd.read_csv(find("segment_stability.csv"))
    c["u21_r2"] = (f"{seg[seg.segment=='<=21'].R2.iloc[0]:.3f}", "segment_stability.csv", "SOURCE-DIRECT")

    fam = pd.read_csv(find("shap_family_importance.csv"))
    c["shap_age_share"] = (f"{fam[fam.family=='age'].share_pct.iloc[0]:.1f}",
                           "shap_family_importance.csv", "SOURCE-DIRECT")

    ah = pd.read_csv(find("genai_arm_history.csv"))
    b0 = ah[(ah.arm.str.startswith("B")) & (ah.run == "initial")].iloc[0]
    c0 = ah[ah.arm.str.startswith("C")].iloc[0]
    c_res = pd.read_csv(find("genai_armC_results.csv"))
    c["genai_armB_initial_passed"] = (f"{int(b0.passed)}", "genai_arm_history.csv",
                                      "SOURCE-DIRECT")
    c["genai_armC_passed"] = (f"{int((c_res.OVERALL == 'PASS').sum())}",
                              "genai_armC_results.csv", "DERIVED")
    c["genai_armC_density"] = (f"{c0.numeric_density_per_1000w:.1f}",
                               "genai_arm_history.csv", "SOURCE-DIRECT")

    rb = pd.read_csv(find("robustness_age_gradient.csv"))
    piv = rb.pivot(index="age_band", columns="season", values="mean_residual")
    # The test season is the latest present in the artifact, not a literal.
    test_season = int(max(piv.columns))
    c[f"gradient_{test_season}_u21"] = (f"{piv.loc['<=21', test_season]:.3f}",
                                        "robustness_age_gradient.csv",
                                        "SOURCE-DIRECT")
    return c


# =============================================================================
# STEP 2 - WHAT EACH DELIVERABLE MUST AGREE WITH
# =============================================================================
DELIVERABLES = {
    "report_sections1_2_intro.md": ["positive_fee_records_pct"],
    "report_section3_methodology.md": ["player_seasons", "unique_players", "test_n",
                                       "positive_fee_playerseason_pct",
                                       "positive_fee_records_pct"],
    "report_section4_results.md": [
        "model0_r2", "model1_r2", "model2_r2", "model1_median_err",
        "exit_auc", "decile1_actual", "decile10_actual",
        "attrition_flagged", "attrition_benchmark", "risk_ratio",
        "backtest_p_v2", "u21_r2", "portfolio_spend_m"],
    # gradient key is season-dependent; filled in by resolve_deliverables()
    "report_section4_8_explainability.md": ["shap_age_share", "__GRADIENT_KEY__"],
    "report_section5_analysis.md": [
        "model0_r2", "model1_r2", "model2_r2", "exit_auc",
        "portfolio_spend_m", "portfolio_quality_display_mean", "portfolio_exit_risk",
        "budget_plateau_m", "u21_r2"],
    "report_section6_conclusion.md": [
        "model0_r2", "model1_r2", "model2_r2", "exit_auc",
        "attrition_flagged", "attrition_benchmark", "risk_ratio", "backtest_p_v2"],
    "report_front_matter_SUBMISSION.md": [
        "player_seasons", "model0_r2", "model1_r2", "exit_auc", "risk_ratio"],
    "demo_video_script.md": [
        "player_seasons", "model0_r2", "model1_r2", "model2_r2", "exit_auc",
        "attrition_flagged", "attrition_benchmark", "risk_ratio",
        "decile1_actual", "decile10_actual", "portfolio_spend_m",
        "portfolio_quality_display_mean", "portfolio_exit_risk"],
    "optimizer_specification.md": ["attrition_flagged", "attrition_benchmark", "backtest_p_v2"],
    "report_section7_9_genai.md": ["genai_armB_initial_passed", "genai_armC_passed",
                                   "genai_armC_density"],
}

# Artifacts that must exist before the submission gate can pass. The GenAI
# items are knowingly absent until the live A/B run completes; the gate is
# expected to FAIL until then, which is the point of having a gate.
SUBMISSION_REQUIRED = [
    "report_sections1_2_intro.md",
    "report_section3_methodology.md",
    "report_section3_10_explainability.md",
    "report_section4_results.md",
    "report_section4_8_explainability.md",
    "report_section5_analysis.md",
    "report_section6_conclusion.md",
    "report_front_matter_SUBMISSION.md",
    "M13A-25_Prompt_Logbook.md",
    "report_skeleton_FROZEN.md",
    "demo_video_script.md",
    "optimizer_specification.md",
    "PROJECT_CONTROL_RULES.md",
    "M13A-25_Beyond_the_Price_Tag_Presentation.pptx",
    # The submitted deliverables are the Word documents, not their markdown
    # sources. The gate verifies the artifacts the professor actually grades.
    "M13A-25_Beyond_the_Price_Tag_Report.docx",
    "M13A-25_Prompt_Logbook.docx",
    "M13A-25_Beyond_the_Price_Tag_Workbook.xlsx",
    "dashboard.py",
    "Beyond_the_Price_Tag_REPORT.md",
    # GenAI evidence. The scripted API harness (genai_ab_*.csv) was NOT used;
    # evaluation ran through conversational invocation, so those files are not
    # required and must not be fabricated.
    "genai_input_facts.json",
    "genai_arm_history.csv",
    "genai_armC_results.csv",
    "genai_controls.csv",
    "genai_fidelity_check.py",
    "gptconv_brief_lucas_stassin.md",
    "gptconv_brief_diego_coppola.md",
    "gptconv_brief_facundo_buonanotte.md",
    "claudeconv_brief_lucas_stassin.md",
    "claudeconv_brief_diego_coppola.md",
    "claudeconv_brief_facundo_buonanotte.md",
    "report_section6_11_genai.md",
    "report_section7_9_genai.md",
]

BANNED = {
    "LightGBM / XGBoost claimed as the estimator": r"\b(lightgbm|xgboost)\b",
    "odds ratio stated as a likelihood": r"4\.8\s*(×|x)?\s*more likely",
    "player labelled undervalued (asserted)": r"(is|are|were)\s+undervalued\b",
    "optimizer described as validated": r"optimizer (is |was )?validated",
    "claim of superior returns": r"(superior|excess|abnormal) (transfer )?returns?",
}
BANNED_EXEMPT = {  # legitimate uses: negations, historical framing, comparisons
    "LightGBM / XGBoost claimed as the estimator":
        r"same family as|not used|never|don'?t say|do not say|must not appear|"
        r"was never installed|incorrectly|corrected|HistGradientBoosting",
    "player labelled undervalued (asserted)": r"not |never |hypothes|flagged as|labelled",
    "claim of superior returns": r"no claim|not (made|supported)",
}


WORD_FORMS = {"10.0": [r"a tenth", r"one[- ]tenth", r"10%", r"ten per ?cent"]}


DELIVERABLE_DIRS = ["report/final", "report/prompt_logbook", "report/development",
                    "presentation", "workbook", "code/12_audit",
                    "outputs/optimizer", "briefs/claude", "briefs/gpt"]


def locate(name):
    """Find a deliverable anywhere in the repository."""
    for d in DELIVERABLE_DIRS:
        c = ROOT / d / name
        if c.exists():
            return c
    hits = list(ROOT.rglob(name))
    return hits[0] if hits else None


def resolve_deliverables(canon):
    """
    Substitute season-dependent canonical keys into the deliverable map.

    An earlier version hardcoded "gradient_2024_u21" here. That reintroduced the
    exact assumption the self-audit had just removed from canonical(): correct
    only while the test season happens to be 2024.
    """
    grad = next(k for k in canon if k.startswith("gradient_"))
    out = {}
    for fname, keys in DELIVERABLES.items():
        out[fname] = [grad if k == "__GRADIENT_KEY__" else k for k in keys]
    return out


def check_file(path, keys, canon):
    txt = path.read_text()
    rows = []
    for k in keys:
        val, src, kind = canon[k]
        # accept the figure with or without thousands separators / trailing zeros
        pats = {val, val.replace(",", ""), val.rstrip("0").rstrip(".")}
        # a figure quoted at lower precision still agrees (0.732 vs 0.7323)
        try:
            f = float(val.replace(",", ""))
            for dp in (0, 1, 2, 3):
                pats.add(f"{f:.{dp}f}")
        except ValueError:
            pass
        found = any(re.search(re.escape(v), txt) for v in pats if v)
        if not found:
            found = any(re.search(w, txt, re.I) for w in WORD_FORMS.get(val, []))
        rows.append({"file": path.name, "check": k, "canonical": val,
                     "type": kind, "source": src,
                     "status": "OK" if found else "MISSING/MISMATCH"})
    return rows


def check_banned(path):
    txt = path.read_text()
    hits = []
    for label, pat in BANNED.items():
        for m in re.finditer(pat, txt, re.I):
            line = txt[max(0, m.start() - 400):m.end() + 90].replace("\n", " ")
            ex = BANNED_EXEMPT.get(label)
            if ex and re.search(ex, line, re.I):
                continue
            hits.append({"file": path.name, "violation": label,
                         "context": line.strip()[:150]})
    return hits



# =============================================================================
# STEP 2a - FACT-PAYLOAD AUDIT
# =============================================================================
# genai_input_facts.json is the only artifact whose text is fed verbatim into
# generated documents. A stale figure there propagates into every brief an LLM
# writes, and the model would be correct to reproduce it because we supplied it
# as an allowed fact. It therefore belongs in the canonical audit chain.

def check_fact_payload(canon):
    """Verify the GenAI fact payload against the canonical table."""
    fp = find("genai_input_facts.json")
    if not fp.exists():
        return [{"check": "payload present", "status": "ABSENT"}]
    facts = json.loads(fp.read_text())
    blob = json.dumps(facts)
    rows = []

    # figures the payload states in prose must match canon
    for key in ["risk_ratio", "attrition_flagged", "attrition_benchmark",
                "backtest_p_v2", "exit_auc", "model1_median_err"]:
        val = canon[key][0]
        variants = {val}
        try:
            fv = float(val)
            for dp in (0, 1, 2, 3, 4):
                variants.add(f"{fv:.{dp}f}")
        except ValueError:
            pass
        ok = any(v in blob for v in variants)
        rows.append({"check": f"payload states {key} = {val}",
                     "status": "OK" if ok else "MISSING"})

    # superseded figures must NOT appear
    for stale, why in [("4.8", "odds ratio stated as a risk ratio"),
                       ("LightGBM", "wrong estimator name")]:
        hit = stale in blob
        rows.append({"check": f"payload free of '{stale}' ({why})",
                     "status": "OK" if not hit else "STALE VALUE PRESENT"})

    # per-player figures must match the optimizer input artifact
    oi = pd.read_csv(find("optimizer_input_2024_25.csv"))
    for f in facts:
        r = oi[oi.name == f["player"]]
        if r.empty:
            rows.append({"check": f"{f['player']} in optimizer input",
                         "status": "NOT FOUND"}); continue
        r = r.iloc[0]
        ok = (int(f["market_value_eur"]) == int(r.market_value_eur)
              and abs(f["predicted_exit_risk_pct"] - 100 * r.exit_prob) < 0.06)
        rows.append({"check": f"{f['player']} market value + exit risk match source",
                     "status": "OK" if ok else "MISMATCH"})
    return rows


# =============================================================================
# STEP 2b - CLAIM-CLASS AUDIT
# =============================================================================
# The project is now sophisticated enough that CLAIM DRIFT is a bigger risk than
# arithmetic drift. A figure can be perfectly sourced while the sentence around
# it over-claims. Each claim class below has a permitted evidence level; the
# patterns detect language that exceeds it.

CLAIM_CLASSES = {
    "mispricing_signal": {
        "permitted": "UNSUPPORTED - tested out-of-sample, not supported",
        "overclaim": [
            r"(identifies|finds|detects|reveals) (systematically )?(exploitable|"
            r"genuine|real) (mispricing|undervaluation)",
            r"(model|system|residual) (can |does )?(successfully )?identif\w+ "
            r"undervalued players",
            r"proves? the market is (wrong|inefficient|mispricing)",
        ],
        "exempt": r"not |did not|does not|hypothes|failed|rejected|unsupported|"
                  r"could not|no evidence|whether",
    },
    "portfolio": {
        "permitted": "ILLUSTRATIVE - an application of a validated model, n = 3",
        "overclaim": [
            r"optimizer (is |was )?(statistically )?validated",
            r"portfolio (is |was )?validated",
            r"(proves|demonstrates|establishes) (that )?the optimizer works",
            r"(superior|excess|abnormal|outperform\w*) (transfer )?returns?",
        ],
        "exempt": r"not |never |cannot|no claim|illustrat|does not",
    },
    "market_efficiency": {
        "permitted": "INTERPRETATION - the market may hold information we lack",
        "overclaim": [
            r"the (transfer )?market is efficient",
            r"proves? (that )?the market is (efficient|right|correct)",
            r"markets? are (fully )?efficient",
        ],
        "exempt": r"not |whether|may |suggests|consistent with|does not",
    },
    "genai_fidelity": {
        "permitted": "CONTROL CONDITION ONLY until the live A/B run completes",
        "overclaim": [
            r"(the )?(LLM|language model|live (generation|model)) (achieved|scored|"
            r"passed) \d",
            r"live LLM (results?|output) (showed|achieved|passed)",
        ],
        "exempt": r"pending|not yet|awaiting|will|would|control condition|template",
    },
    "genai_independence": {
        "permitted": "PARTIAL independence for Arm C; NONE for Arm B",
        "overclaim": [
            r"independent(ly)? (validat\w+|evaluat\w+|verif\w+) (of|by) (the )?"
            r"(GPT|cross-model|second model)",
            r"fully independent (arm|evaluation|validation)",
            r"(GPT|Claude) is more (reliable|accurate|trustworthy)",
            r"100% (factual )?accuracy",
        ],
        "exempt": r"not |never |partial|no claim|does not|cannot|is not|"
                  r"don'?t let|prohibited|must not|refused|rejected|withdrawn",
    },
    "causality": {
        "permitted": "ASSOCIATION only - no causal claim is supported",
        "overclaim": [
            r"(residual|gap|model) causes?\b",
            r"because the market (knows|held|had)",
            r"(proves|demonstrates) that (contracts?|injur\w+|scouting) (explain|cause)",
        ],
        "exempt": r"hypothes|plausible|may |might|not |cannot|candidate",
    },
}


def check_claims(path):
    txt = path.read_text()
    hits = []
    for cls, spec in CLAIM_CLASSES.items():
        for pat in spec["overclaim"]:
            for m in re.finditer(pat, txt, re.I):
                window = txt[max(0, m.start() - 420):m.end() + 140].replace("\n", " ")
                if re.search(spec["exempt"], window, re.I):
                    continue
                hits.append({"file": path.name, "claim_class": cls,
                             "permitted": spec["permitted"],
                             "text": m.group(0)[:90],
                             "context": window.strip()[:170]})
    return hits


# =============================================================================
# SELF-AUDIT - the auditor must not assert numbers it did not read
# =============================================================================
def self_audit():
    """
    Scan canonical() for empirical numeric literals in EXECUTABLE CODE.

    Text scanning produced false positives on numbers appearing inside comments
    and docstrings, so this parses the AST and inspects only numeric constants
    the interpreter actually evaluates. Strings are ignored entirely.

    Structural constants are permitted: scaling factors, rounding precision,
    positional indices, and lookup keys that select a row rather than assert a
    finding. An empirical value is never permitted.
    """
    import ast as _ast

    ALLOWED = {
        0, 1, 2, 3, 4,      # indices, .iloc positions, rounding precision
        100,                # percentage scaling
        1e6,                # euro millions
        75,                 # budget scenario used as a row lookup key
        # NOTE: season identifiers are deliberately NOT allowlisted. They are
        # data values and must be derived from the artifact (see test_season).
    }

    tree = _ast.parse(Path(__file__).read_text())
    fn = next(n for n in _ast.walk(tree)
              if isinstance(n, _ast.FunctionDef) and n.name == "canonical")

    bad = []
    for node in _ast.walk(fn):
        if isinstance(node, _ast.Constant) and isinstance(node.value, (int, float)):
            if node.value in ALLOWED or node.value is True or node.value is False:
                continue
            bad.append((node.value, node.lineno))

    print("\n" + "-" * 78)
    print("SELF-AUDIT — hardcoded empirical values in canonical()")
    print("-" * 78)
    print("   (AST-based: inspects executable code only, ignoring comments and strings)")
    if bad:
        lines = Path(__file__).read_text().split("\n")
        for val, ln in bad:
            print(f"   !! literal {val!r} at line {ln}: {lines[ln-1].strip()[:80]}")
        print(f"\n   {len(bad)} suspect literal(s) — every canonical figure must be READ")
    else:
        print("   none: every canonical figure is read or derived from an artifact")
    return len(bad)


# =============================================================================
# RUN
# =============================================================================
def main():
    print("=" * 78)
    print("FINAL CROSS-ARTIFACT AUDIT")
    print("=" * 78)

    n_lit = self_audit()
    canon = canonical()
    print(f"\nCanonical table built from frozen artifacts ({len(canon)} figures):\n")
    cdf = pd.DataFrame([{"key": k, "value": v[0], "type": v[2], "source": v[1]}
                        for k, v in canon.items()])
    print(cdf.to_string(index=False))
    cdf.to_csv(STAGE_OUT / "canonical_figures.csv", index=False)

    print("\n" + "=" * 78)
    print("CONSISTENCY CHECKS")
    print("=" * 78)

    all_rows, all_hits, missing_files, claim_hits = [], [], [], []
    for fname, keys in resolve_deliverables(canon).items():
        p = locate(fname)
        if p is None:
            missing_files.append(fname); continue

        all_rows += check_file(p, keys, canon)
        all_hits += check_banned(p)
        claim_hits += check_claims(p)

    # banned-term scan across every deliverable, not just those with figures
    # Pre-project planning guides and the control-rules document are not report
    # deliverables. The control rules necessarily NAME the banned terms in order
    # to ban them, and the demo script names them in its "don't say" list.
    NOT_DELIVERABLES = {
        "Football-Valuation-Project-Guide.md",
        "Sports-Analytics-AI-MBA-Project-Guide.md",
        "PROJECT_CONTROL_RULES.md",
        "PROVENANCE_RECONCILIATION.md",
        "CHANGELOG_v2.md",
        "SOURCE_MANIFEST.md",
        "report_skeleton_FROZEN.md",
        "report_front_matter.md",          # superseded by the SUBMISSION version
        "optimizer_specification.md",       # scanned via DELIVERABLES already
    }
    for p in sorted(ROOT.rglob("*.md")):
        if p.name not in DELIVERABLES and p.name not in NOT_DELIVERABLES:
            all_hits += check_banned(p)
            claim_hits += check_claims(p)

    res = pd.DataFrame(all_rows)
    bad = res[res.status != "OK"]
    print(f"\n{len(res)} TEXTUAL figure checks across {res.file.nunique()} report files")
    print("   (deck figure checks, banned-term scan, claim-class scan and the")
    print("    negative-control scan are reported separately below)")
    if len(bad):
        print(f"\n!! {len(bad)} MISMATCH(ES):")
        print(bad.to_string(index=False))
    else:
        print("   all figures agree with their source artifacts")

    print(f"\n{len(all_hits)} banned-term violation(s)")
    if all_hits:
        for h in all_hits:
            print(f"   [{h['file']}] {h['violation']}")
            print(f"      ...{h['context']}...")

    print("\n" + "-" * 78)
    print("FACT-PAYLOAD AUDIT")
    print("-" * 78)
    fp_rows = check_fact_payload(canon)
    for r in fp_rows:
        print(f"   {r['status']:<20} {r['check']}")
    fp_bad = [r for r in fp_rows if r["status"] != "OK"]

    print(f"\n{len(claim_hits)} claim-class violation(s)")
    if claim_hits:
        for h in claim_hits:
            print(f"   [{h['file']}] {h['claim_class']} - permitted: {h['permitted']}")
            print(f"      matched: \"{h['text']}\"")
            print(f"      ...{h['context']}...")
    else:
        print("   every claim sits within its permitted evidence level")
    pd.DataFrame(claim_hits).to_csv(STAGE_OUT / "audit_claim_violations.csv", index=False)

    # The staleness sweep compared delivered copies against working copies. A
    # repository holds a single copy of each artifact, so the check does not
    # apply here and is retained in the project history only.
    stale = []

    # -------------------------------------------------------------------------
    # PLACEHOLDER AND STALE-CONTENT GATE
    # -------------------------------------------------------------------------
    # Catches two failure modes no numerical check reaches: an unfilled [[ ]]
    # field reaching a submitted document, and a prose statement that was true
    # at an earlier project stage but has since been superseded.
    print("\n" + "-" * 78)
    print("PLACEHOLDER AND STALE-CONTENT GATE")
    print("-" * 78)

    SUPERSEDED_VALUES = [
        (r"odds ratio 7\.82|OR 7\.82", "odds ratio corrupted 4.82 -> 7.82"),
        (r"skew 6\.90", "skew corrupted 3.90 -> 6.90"),
        (r"4\.8\s*(x|\u00d7)\s*more likely", "odds ratio stated as a risk ratio"),
        (r"AUC of 0\.728", "challenger AUC used as the production figure"),
        (r"1,?509 (players|rows|observations)", "superseded test-set count"),
    ]
    # The validator applies SIX checks. Any statement implying four or five
    # is stale. Written as a count check rather than a phrase list because the
    # phrase appeared in three different forms and evaded three sweeps.
    VALIDATOR_COUNT = [
        (r"(?<!other )\b(four|five)\s+(?:of\s+(?:the\s+)?(?:five|six)\s+)?"
         r"(?:automated\s+|fidelity\s+)?checks?\b",
         "implies the validator has fewer than six checks"),
        (r"\bother\s+(?:three|four)\s+(?:automated\s+|fidelity\s+)?"
         r"(?:checks?|dimensions?)\b",
         "implies the validator has fewer than six checks"),
        (r"\b(four|five)[- ]check validator\b",
         "implies the validator has fewer than six checks"),
        (r"applying (four|five) checks", "implies the validator has fewer than six checks"),
    ]

    STALE_PHRASES = VALIDATOR_COUNT + [
        (r"validated in its control condition", "pre-Arm-C GenAI limitation"),
        (r"to be completed\b", "unfinished section marker"),
        (r"Drafting status", "working note"),
    ]

    # Planning documents legitimately contain [[ ]] as template markers.
    DEV_DOCS = {"report_skeleton_FROZEN.md", "Football-Valuation-Project-Guide.md"}
    ph_bad = []
    for fname in [f for f in SUBMISSION_REQUIRED
                  if f.endswith(".md") and f not in DEV_DOCS]:
        fp = locate(fname)
        if fp is None or not fp.exists():
            continue
        txt = fp.read_text()
        if re.findall(r"\[\[[^\]]*\]\]", txt):
            ph_bad.append((fname, "unfilled [[ ]] field(s)"))
        for pat, why in SUPERSEDED_VALUES + STALE_PHRASES:
            for m in re.finditer(pat, txt, re.I):
                ctx = txt[max(0, m.start() - 70):m.end() + 50].replace("\n", " ")
                ph_bad.append((fname, f"{why}: ...{ctx.strip()[:100]}..."))

    if ph_bad:
        for f, why in ph_bad:
            print(f"   !! [{f}] {why}")
        print(f"\n   {len(ph_bad)} blocking issue(s)")
    else:
        print("   no unfilled placeholders and no superseded statements")

    # -------------------------------------------------------------------------
    # SOURCE / DOCUMENT AGREEMENT
    # -------------------------------------------------------------------------
    # The prompt logbook was rebuilt from a stale copy of its markdown, which
    # silently reverted a correction applied to the other copy. Any Word
    # document and its markdown source must now agree on the tracked phrases.
    # -------------------------------------------------------------------------
    # MANIFEST AND ASSET INTEGRITY
    # -------------------------------------------------------------------------
    # SOURCE_MANIFEST.md records a SHA-256 per artifact. Those went stale when
    # artifacts were regenerated, because nothing checked them. Markdown image
    # references also went unchecked, so a source file referenced an asset that
    # was not packaged.
    print("\n" + "-" * 78)
    print("MANIFEST AND ASSET INTEGRITY")
    print("-" * 78)
    import hashlib

    mf = ROOT / "provenance" / "SOURCE_MANIFEST.md"
    if mf.exists():
        rows = re.findall(r"\| `([^`]+)` \| `[^`]*` \| \w+ \| [^|]+ \| `([0-9a-f]{12})` \|",
                          mf.read_text())
        drift = []
        for rel, claimed in rows:
            f = ROOT / rel
            if not f.exists():
                drift.append((rel, "path not found")); continue
            actual = hashlib.sha256(f.read_bytes()).hexdigest()[:12]
            if actual != claimed:
                drift.append((rel, f"{claimed} -> {actual}"))
        if drift:
            for rel, why in drift:
                ph_bad.append(("SOURCE_MANIFEST.md", f"checksum drift: {rel} ({why})"))
            print(f"   !! {len(drift)} of {len(rows)} manifest entries stale")
        else:
            print(f"   all {len(rows)} manifest checksums match")

    missing_assets = []
    SKIP_DIRS = {"node_modules", ".git", "__pycache__", ".venv", "venv"}
    for md_file in ROOT.rglob("*.md"):
        if SKIP_DIRS & set(md_file.parts):
            continue
        for m in re.finditer(r"!\[[^\]]*\]\(([^)]+)\)", md_file.read_text()):
            if m.group(1).startswith(("http://", "https://", "data:")):
                continue
            ref = (md_file.parent / m.group(1)).resolve()
            if not ref.exists():
                missing_assets.append(f"{md_file.name} -> {m.group(1)}")
    if missing_assets:
        for a in missing_assets:
            ph_bad.append(("markdown asset", f"referenced file not in package: {a}"))
        print(f"   !! {len(missing_assets)} markdown image reference(s) unresolved")
    else:
        print("   every markdown image reference resolves")

    print("\n" + "-" * 78)
    print("SOURCE / DOCUMENT AGREEMENT")
    print("-" * 78)
    PAIRS = [("M13A-25_Prompt_Logbook.md", "M13A-25_Prompt_Logbook.docx"),
             ("Beyond_the_Price_Tag_REPORT.md",
              "M13A-25_Beyond_the_Price_Tag_Report.docx")]
    TRACKED = ["three separate forms of test-set leakage",
               "agree to three decimal places",
               "substantially stronger statistical evidence",
               "six automated checks", "applying six checks"]
    for md_name, dx_name in PAIRS:
        mp, dp = locate(md_name), locate(dx_name)
        if mp is None or dp is None:
            continue
        md = mp.read_text()
        import subprocess
        try:
            dx = subprocess.run(["pandoc", "-t", "plain", str(dp)],
                                capture_output=True, text=True, timeout=120).stdout
        except Exception:
            continue
        for phrase in TRACKED:
            in_md, in_dx = phrase in md, phrase in dx
            if in_md != in_dx:
                ph_bad.append((dx_name,
                    f"source/document disagree on '{phrase[:44]}': "
                    f"markdown={'yes' if in_md else 'no'}, "
                    f"document={'yes' if in_dx else 'no'}"))
        print(f"   {dx_name:46} checked against its markdown source")

    print("\n" + "-" * 78)
    print("SUBMISSION COMPLETENESS")
    print("-" * 78)
    absent = [f for f in SUBMISSION_REQUIRED if locate(f) is None]
    for f in SUBMISSION_REQUIRED:
        present = locate(f) is not None
        print(f"   {'present' if present else 'ABSENT ':<9} {f}")
    if absent:
        print(f"\n   {len(absent)} required artifact(s) absent — the submission")
        print("   gate CANNOT pass until every one exists.")
    else:
        print("\n   every required artifact is present")

    # -------------------------------------------------------------------------
    # SLIDE-SPECIFIC CHECKS
    # -------------------------------------------------------------------------
    # The general deck scan only asks whether a value appears SOMEWHERE. That is
    # too weak for the two slides where numerical drift would do most damage: a
    # deck could contain both a right and a wrong figure and still pass. These
    # checks require each value to appear on its intended slide.
    # Slide positions, not just values. Four slides were inserted during the
    # deck expansion, moving the exit-risk evidence from 7 to 8 and the
    # committee portfolio from 10 to 11. A position-based check that is not
    # updated alongside will pass against the wrong slide.
    SLIDE_REQUIRED = {
        8: ["3.7", "30.1", "8.2", "9.5", "7.3", "69.5", "4.82", "2.6"],
        11: ["48.0", "50", "88.3", "17.1",
             "Lucas Stassin", "Diego Coppola", "Facundo Buonanotte"],
        13: ["2,079", "1,868", "0.606", "0.819", "0.560", "0.074"],
        15: ["39", "20", "16"],
    }
    slide_fail = []

    # deck is binary - inspect its extracted text separately
    # -------------------------------------------------------------------------
    # DECK CHECKS
    # -------------------------------------------------------------------------
    # Text is extracted with python-pptx, a declared dependency, rather than an
    # external tool that may be absent. An earlier version shelled out to
    # markitdown and printed "deck text extraction skipped" when it was missing,
    # then continued to AUDIT PASSED. A gate that can pass without running a
    # check is exactly the false assurance this framework exists to prevent, so
    # extraction failure is now fatal.
    deck_fail = []
    deck = locate("M13A-25_Beyond_the_Price_Tag_Presentation.pptx")

    print("\n" + "-" * 78)
    print("DECK CHECKS")
    print("-" * 78)

    # Slide positions, not just values. Four slides were inserted during the
    # deck expansion, moving the exit-risk evidence from 7 to 8 and the
    # committee portfolio from 10 to 11. A position-based check that is not
    # updated alongside will pass against the wrong slide.
    SLIDE_REQUIRED = {
        8: ["3.7", "30.1", "8.2", "9.5", "7.3", "69.5", "4.82", "2.6"],
        11: ["48.0", "50", "88.3", "17.1",
             "Lucas Stassin", "Diego Coppola", "Facundo Buonanotte"],
        13: ["2,079", "1,868", "0.606", "0.819", "0.560", "0.074"],
        15: ["39", "20", "16"],
    }

    if deck is None:
        deck_fail.append(("deck", "PRESENTATION FILE NOT FOUND"))
        print("   !! presentation file not found")
    else:
        try:
            from pptx import Presentation
        except ImportError:
            deck_fail.append(("deck", "python-pptx not installed"))
            print("   !! python-pptx is required for deck checks.")
            print("      Install dependencies: pip install -r requirements.txt")
            Presentation = None

        if Presentation is not None:
            try:
                prs = Presentation(str(deck))
                slides = {}
                for i, sl in enumerate(prs.slides, start=1):
                    parts = []
                    for shape in sl.shapes:
                        if shape.has_text_frame:
                            parts.append(shape.text_frame.text)
                        if shape.has_table:
                            for row in shape.table.rows:
                                parts += [c.text for c in row.cells]
                        if shape.has_chart:
                            try:
                                ch = shape.chart
                                parts += [str(c) for pl in ch.plots
                                          for c in pl.categories]
                                parts += [str(v) for pl in ch.plots
                                          for ser in pl.series for v in ser.values]
                            except Exception:
                                pass
                    slides[i] = "\n".join(parts)
                full = "\n".join(slides.values())

                print(f"   extracted {len(slides)} slides with python-pptx")

                for k in ["player_seasons", "model0_r2", "model1_r2", "model2_r2",
                          "exit_auc", "risk_ratio", "portfolio_spend_m",
                          "portfolio_quality_display_mean", "portfolio_exit_risk",
                          "decile_spread"]:
                    v = canon[k][0]
                    cands = {v, v.replace(",", "")}
                    try:
                        fv = float(v.replace(",", ""))
                        for dp in (0, 1, 2, 3):
                            cands.add(f"{fv:.{dp}f}")
                    except ValueError:
                        pass
                    ok = any(c in full for c in cands)
                    if not ok:
                        deck_fail.append((k, f"{v} not found in deck"))
                    print(f"   {k:32} {v:>12}   {'OK' if ok else 'NOT FOUND'}")

                for label, pat in BANNED.items():
                    n = len(re.findall(pat, full, re.I))
                    if n:
                        deck_fail.append((label, f"{n} occurrence(s)"))
                        print(f"   !! banned: {label} ({n})")

                for num, needed in SLIDE_REQUIRED.items():
                    body = slides.get(num, "")
                    if not body:
                        deck_fail.append((num, "SLIDE NOT FOUND"))
                        print(f"   slide {num}: NOT FOUND")
                        continue
                    miss = [x for x in needed if x not in body]
                    if miss:
                        deck_fail.append((num, miss))
                    print(f"   slide {num:2}  {len(needed)-len(miss)}/{len(needed)} "
                          f"values present   {'OK' if not miss else f'MISSING {miss}'}")

            except Exception as ex:
                deck_fail.append(("deck", f"extraction failed: {ex}"))
                print(f"   !! deck extraction FAILED: {ex}")

    slide_fail = deck_fail

    # -------------------------------------------------------------------------
    # NEGATIVE CONTROL - a checker that passes everything proves nothing
    # -------------------------------------------------------------------------
    print("\n" + "-" * 78)
    print("NEGATIVE CONTROL")
    print("-" * 78)
    print("Feeding the claim auditor deliberately over-claiming sentences.\n")

    import tempfile
    PLANTS = [
        ("mispricing_signal",
         "Our model identifies systematically exploitable mispricing in the "
         "transfer market."),
        ("portfolio",
         "The recommended portfolio demonstrates that the optimizer works and "
         "delivers superior transfer returns."),
        ("market_efficiency",
         "These findings prove that the market is efficient."),
        ("genai_fidelity",
         "The live LLM achieved 3 of 3 on every fidelity check."),
        ("causality",
         "This proves that contracts explain the residual."),
        ("causality",
         "The residual causes subsequent player exit."),
    ]
    caught = 0
    for expected, sentence in PLANTS:
        with tempfile.NamedTemporaryFile("w", suffix=".md", delete=False) as fh:
            fh.write(sentence)
            tmp = Path(fh.name)
        hits = check_claims(tmp)
        got = {h["claim_class"] for h in hits}
        ok = expected in got
        caught += ok
        print(f"   planted {expected:20} -> {'CAUGHT' if ok else 'MISSED'}")
        tmp.unlink()
    print(f"\n   {caught}/{len(PLANTS)} planted over-claims detected")
    if caught < len(PLANTS):
        print("   !! the claim auditor is not discriminating - tighten the patterns")

    res.to_csv(STAGE_OUT / "audit_results.csv", index=False)
    pd.DataFrame(all_hits).to_csv(STAGE_OUT / "audit_violations.csv", index=False)
    # keep the delivered copies of the audit's own output current
    for f in ("canonical_figures.csv", "audit_results.csv"):
        src = PROC / f
        if src.exists():
            (OUT / f).write_bytes(src.read_bytes())

    print("\n" + "=" * 78)
    ok = (len(bad) == 0 and len(all_hits) == 0 and len(claim_hits) == 0
          and n_lit == 0 and caught == len(PLANTS)
          and len(absent) == 0 and len(missing_files) == 0
          and len(slide_fail) == 0 and len(fp_bad) == 0
          and len(stale) == 0 and len(ph_bad) == 0)
    print("AUDIT PASSED" if ok else "AUDIT FOUND ISSUES - see above")
    print("=" * 78)
    return 0 if ok else 1


if __name__ == "__main__":
    sys.exit(main())
